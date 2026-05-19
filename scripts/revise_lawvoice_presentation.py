from pathlib import Path
from textwrap import dedent

from pptx import Presentation


INPUT_PPTX = Path(r"C:\Users\xsanf\Downloads\Telegram Desktop\LawVoice_Defense_Presentation_EN.pptx")
OUTPUT_DIR = Path(r"F:\GitHub\vangZ_strict_patched_plus_voice_assistant\docs\defense")
OUTPUT_PPTX = OUTPUT_DIR / "LawVoice_Defense_Presentation_EN_revised.pptx"
OUTPUT_SPEECH = OUTPUT_DIR / "LawVoice_Defense_Speech_EN_8min.md"


SLIDE_TEXT = {
    1: {
        6: "Current MVP: allaw-urist.ru",
    },
    2: {
        14: "LawVoice is a voice-first educational assistant for helping teenagers explore legal situations.",
        16: "Current state",
        17: "Current MVP interface, LawVoice personas, dialogue logic, knowledge base, voice flow, and action-plan prototype.",
        19: "Still in progress",
        20: "Student name, supervisor, pilot-study results, evaluation metrics, and the final expanded bibliography.",
    },
    3: {
        10: "This work focuses on the design and current MVP implementation of LawVoice, a voice-based educational assistant that helps teenagers discuss everyday legal situations in a safe format.",
        11: "- The project is positioned as an educational voice simulator, not as a legal consultation service.\n- The current MVP centers on cyberbullying, online purchases, and rights at school.\n- The implemented dialogue flow clarifies facts, explains risks, suggests options, and ends with a concrete next step.\n- The LawVoice layer is built on top of an adapted OpenAI realtime voice-agent base.",
        14: "The aim is not simply to give an answer, but to guide the user through a structured chain: facts -> risks -> options -> next step.",
        17: "Voice-first - educational - scenario-based",
    },
    4: {
        10: "- LawVoice addresses everyday situations that teenagers may face outside a formal law class.\n- Voice interaction lowers the entry barrier because speaking through a case can be easier than reading long legal texts.\n- The educational value is connected not only with legal literacy, but also with critical thinking, because the user compares options and consequences.\n- The project combines AI dialogue, safe communication rules, and a practical learning interface in one MVP.",
        14: "Cyberbullying\n\nOnline purchases\n\nRights at school",
    },
    5: {
        13: "To develop and present the current MVP of a voice-based intelligent assistant that helps teenagers explore legal situations and practice structured reasoning.",
        17: "1. Analyze teen legal scenarios and define a safe educational framing.\n\n2. Design dialogue logic, personas, age adaptation, and escalation rules.\n\n3. Implement the current MVP interface, voice interaction, and knowledge-base-supported action planning.\n\n4. Define limitations and next steps for legal, pedagogical, and empirical evaluation.",
    },
    6: {
        29: "At the current project stage, steps 1-4 are represented in the MVP. Steps 5-6 remain the next research and deployment stage.",
    },
    7: {
        12: "The teenager selects a situation and one of the LawVoice personas.",
        15: "The assistant asks short clarifying questions without requesting unnecessary personal data.",
        18: "LawVoice explains rights, obligations, and risks in plain language.",
        21: "The user receives 2-3 possible legal and safe next steps.",
        24: "The dialogue ends with a recommended action and a short 24-hour checklist.",
        25: "Educational outcome: the user practices structured reasoning instead of waiting for a ready-made verdict.",
    },
    8: {
        11: "Persona profiles",
        12: "Mentor, peer, analyst, and safety personas with voice and style settings.",
        14: "Adaptive mode",
        15: "Teen and 18+ modes with runtime adjustment for risk and anxiety.",
        18: "LawVoice system prompt and dialogue rules at session start.",
        20: "Scenario and materials",
        21: "Scenario hints, supporting materials, and prepared context for the conversation.",
        24: "Document upload, chunking, search, and token estimates.",
        27: "Realtime session, microphone, audio visualization, and context prompts.",
        30: "Grounded draft plan with steps, priorities, and comments.",
        33: "Session log, message count, token usage, and API cost.",
    },
    9: {
        13: "- profile selection and dialogue-style screen\n- age mode and prompt editing\n- scenario and document input\n- voice-session controls\n- audio, conversation, and log widgets",
        17: "- intent routing and LawVoice dialog manager\n- session state and safety-aware profile switching\n- grounded action-plan generation\n- analytics logging and prompt control",
        21: "- OpenAI Realtime API and WebRTC\n- Node.js / Express backend\n- PostgreSQL with pgvector and knowledge chunks\n- gpt-realtime-mini and embedding services",
        23: "Current implementation: a LawVoice domain layer built on top of an adapted realtime voice-agent foundation.",
    },
    10: {
        13: "- clarify what happened, where it happened, and whether evidence exists\n- separate facts from emotion and assess risk\n- suggest safe next steps and evidence preservation\n- in dangerous situations, recommend a trusted adult or 112",
        17: "- clarify the order, payment, delivery, and seller\n- review evidence such as receipts, chats, screenshots, and order cards\n- discuss return, complaint, and documentation of the problem\n- offer a practical next step for today",
        21: "- discuss the conflict calmly and without stigma\n- explain rights, obligations, and school boundaries in simple language\n- help formulate a message for a teacher, parent, or administrator\n- turn the situation into a short action plan",
    },
    11: {
        12: "Profile selection\nAge mode\nRules and instructions\nScenario and document input",
        15: "Start and end session\nMicrophone\nAudio input and output\nUrgent and quiet prompts",
        18: "Draft action plan\nLearning analytics\nAction log\nTechnical panel and cost",
        20: "This layout separates preparation, live dialogue, and post-session reflection, which is useful both for a demo and for later research.",
    },
    12: {
        10: "- The current version is an MVP, and some interface and logic blocks are still demonstrational rather than research-validated.\n- There are no completed empirical results or classroom pilot measurements yet.\n- The legal content is currently tailored to Russian-language scenarios and still needs expert legal and pedagogical review.\n- Educational deployment would require clearer privacy, moderation, and escalation procedures.",
        13: "The presentation intentionally avoids unsupported claims about learning impact or legal accuracy.",
    },
    13: {
        12: "Expansion of scenarios, documents, explanation modes, and role profiles.",
        15: "Refinement of escalation rules, moderation logic, and privacy boundaries.",
        18: "Pilot testing and measurement of legal-literacy and critical-thinking outcomes.",
        19: "- Add new scenarios, including peer-pressure and detention-related support.\n- Replace remaining generic elements with fully LawVoice-specific components.\n- Prepare pilot testing in an educational environment.\n- Develop metrics for response quality, safety, and educational usefulness.",
    },
    14: {
        11: "1. The project defines a feasible concept of a voice-first educational legal assistant for teenagers.\n\n2. The current MVP shows that realtime voice interaction can support structured reasoning: facts, risks, options, and the next step.\n\n3. The prototype provides a basis for future expert review, pilot testing, and impact evaluation rather than finalized legal guidance.",
    },
    15: {
        10: "1. LawVoice MVP. https://allaw-urist.ru/\n\n2. OpenAI. Realtime API documentation. https://platform.openai.com/docs/api-reference/realtime-sessions\n\n3. LawVoice project repository. README and current MVP implementation.\n\n4. Project note: LawVoice General Character Prompt.\n\n5. Project note: Voice Assistant Architecture.",
    },
    16: {
        14: "Persona and dialogue style\nAge mode\nRules\nBase prompt\nProfiles",
        18: "Voice dialogue\nStart / end session\nMicrophone / audio visualization",
        22: "Dialogue\nAction log\nDraft action plan",
        26: "Scenario hints and materials\nKnowledge base",
        30: "Learning analytics\nTechnical panel\nAPI cost",
    },
    17: {
        12: "The user describes the case or chooses a scenario",
        16: "The assistant clarifies key details",
        20: "LawVoice explains rights and safety factors in simple language",
        24: "2-3 possible next steps",
        28: "Recommended action + 24-hour checklist",
        30: "This flow is suitable both for a product demonstration and for a later educational study.",
    },
}


NOTES = {
    1: dedent(
        """\
        Good afternoon. My thesis presentation is about LawVoice, a voice-based intelligent assistant for teenagers. The current MVP is available as a web prototype and is designed as an educational tool for discussing everyday legal situations and practicing structured reasoning.
        """
    ).strip(),
    2: dedent(
        """\
        I will briefly cover the subject area, the relevance of the topic, the goal of the work, the current MVP, its architecture, the main user scenarios, the current limitations, and the next development steps. I will focus on what is already implemented and clearly separate it from what is still planned.
        """
    ).strip(),
    3: dedent(
        """\
        The subject area of this work is a voice-first assistant that helps teenagers talk through legal situations in a safe educational format. LawVoice is not positioned as a replacement for a lawyer. Instead, it structures the dialogue: first facts, then risks, then options, and finally a concrete next step. In technical terms, the current LawVoice prototype is a domain-specific layer built on top of an adapted realtime voice-agent base.
        """
    ).strip(),
    4: dedent(
        """\
        The topic is relevant for several reasons. Teenagers often face legal or quasi-legal conflicts outside a formal classroom, but they may not be ready to read long legal explanations. A voice format lowers the barrier to entry and makes the interaction more natural. At the same time, the value of the system is not only legal literacy but also critical thinking, because the user is encouraged to compare options and consequences instead of searching for one short answer.
        """
    ).strip(),
    5: dedent(
        """\
        The goal of the work is to develop and present the current MVP of a voice-based assistant that helps teenagers explore legal situations and practice structured reasoning. To reach this goal, I define the relevant scenarios, design safe dialogue logic and personas, implement the interface and voice interaction, and then identify the limits of the prototype together with the criteria for future evaluation.
        """
    ).strip(),
    6: dedent(
        """\
        This roadmap shows the logic of the project. In the current state, the first four stages are already represented at MVP level: subject analysis, scenario and role design, dialogue logic, and interface implementation. The last two stages, which are pilot testing and impact evaluation, are intentionally treated as future work because they require a separate empirical study and formal measurement design.
        """
    ).strip(),
    7: dedent(
        """\
        The user journey is intentionally simple. First, the teenager chooses a scenario and a LawVoice persona. Then the assistant asks for the minimum facts needed to understand the case without collecting unnecessary personal data. After that, LawVoice explains the relevant rights, obligations, and risks in plain language, offers several possible next steps, and finishes with one recommended action plus a short checklist for the next 24 hours. In this way, the conversation supports reasoning rather than passive advice consumption.
        """
    ).strip(),
    8: dedent(
        """\
        This slide summarizes what is already present in the current MVP. The interface includes several LawVoice personas, including mentor, peer, analyst, and safety roles. It supports teen and adult modes, a base system prompt, scenario and material input, a knowledge base for uploaded documents, realtime voice dialogue, draft action-plan generation, and analytics such as message count, token usage, and API cost. So the prototype is already more than a static concept slide deck; it is an interactive working MVP.
        """
    ).strip(),
    9: dedent(
        """\
        The architecture can be described in three layers. The interface layer contains profile selection, age mode, prompt editing, document input, voice controls, and visual widgets. The logic layer contains intent routing, the LawVoice dialog manager, safety-aware profile switching, and grounded action-plan generation. The AI and data layer combines the OpenAI Realtime API, WebRTC, a Node.js and Express backend, PostgreSQL with pgvector for knowledge chunks, and the current gpt-realtime-mini model. This is why I describe the project as a LawVoice domain layer built on top of an adapted realtime infrastructure.
        """
    ).strip(),
    10: dedent(
        """\
        At the moment, the MVP focuses on three main scenario groups. The first is cyberbullying, where the system helps clarify facts, separate evidence from emotion, and suggest safe next actions. The second is online purchases, where the assistant helps structure the situation around payment, delivery, receipts, screenshots, and possible complaint or return steps. The third is rights at school, where the goal is to explain obligations and boundaries in simple language and help the teenager prepare a calm and practical next step. Across all three scenarios, safety escalation remains a cross-cutting rule.
        """
    ).strip(),
    11: dedent(
        """\
        The interface is organized into three functional blocks. First, there is the preparation block, which includes persona selection, age mode, rules, prompt configuration, and document input. Second, there is the live dialogue block with the voice session, microphone, audio feedback, and context prompts. Third, there is the reflection and control block with the draft action plan, analytics, logs, and technical monitoring. This structure is convenient both for a demo session and for a future educational study.
        """
    ).strip(),
    12: dedent(
        """\
        It is important to state the limitations clearly. This version is still an MVP, so some interface and logic blocks are demonstrational rather than research-validated. There is no completed classroom pilot yet, so I do not claim measured educational impact. The current legal content is tailored to Russian-language scenarios and still needs expert legal and pedagogical review. In addition, any real educational deployment would require stronger privacy, moderation, and escalation procedures. So this presentation stays intentionally conservative in its claims.
        """
    ).strip(),
    13: dedent(
        """\
        The next development stage has three directions. The first is content expansion: more scenarios, documents, explanation modes, and role profiles. The second is safety and product maturity: stronger moderation logic, clearer escalation rules, and the replacement of remaining generic elements with fully LawVoice-specific components. The third is evaluation: a pilot study in an educational environment together with metrics for response quality, safety, legal literacy, and critical-thinking support.
        """
    ).strip(),
    14: dedent(
        """\
        In conclusion, the current result is a feasible concept and MVP of a voice-first educational legal assistant for teenagers. The prototype demonstrates that realtime voice interaction can support a reasoning process built around facts, risks, options, and the next step. At the same time, the project should be understood as a basis for further expert review and empirical testing, not as a finalized legal guidance product.
        """
    ).strip(),
    15: dedent(
        """\
        The presentation is based on three types of sources: the deployed LawVoice MVP, official OpenAI documentation for the realtime stack, and the current project documentation and implementation materials. I separate these sources deliberately so that the defense remains tied to the real current state of the project.
        """
    ).strip(),
    16: dedent(
        """\
        This is a backup slide. If needed, I can use it to show how the main screen is divided into preparation, live voice dialogue, and post-session reflection blocks.
        """
    ).strip(),
    17: dedent(
        """\
        This is also a backup slide. It summarizes the educational session logic as a short chain from situation to facts, rules and risks, options, and finally a recommended action with a 24-hour checklist.
        """
    ).strip(),
    18: dedent(
        """\
        Thank you for your attention. I will be glad to answer questions about the dialogue logic, the architecture, or the next research stage.
        """
    ).strip(),
}


TIMINGS = {
    1: "~0:20",
    2: "~0:20",
    3: "~0:35",
    4: "~0:40",
    5: "~0:35",
    6: "~0:35",
    7: "~0:45",
    8: "~0:45",
    9: "~0:50",
    10: "~0:50",
    11: "~0:40",
    12: "~0:45",
    13: "~0:40",
    14: "~0:35",
    15: "~0:20",
    16: "backup",
    17: "backup",
    18: "~0:10",
}


def set_shape_text(slide, shape_index, text):
    shape = slide.shapes[shape_index]
    if not getattr(shape, "has_text_frame", False):
        raise RuntimeError(f"Slide shape {shape_index} has no text frame")
    shape.text = text


def set_notes(slide, text):
    notes_frame = slide.notes_slide.notes_text_frame
    notes_frame.text = text


def build_speech_markdown():
    lines = [
        "# LawVoice Defense Speech (8-Minute Version)",
        "",
        "Main flow: slides 1-15 and 18.",
        "Slides 16-17 are backup slides for questions or extra time.",
        "",
    ]
    for slide_number in range(1, 19):
        timing = TIMINGS.get(slide_number, "")
        label = f"## Slide {slide_number}"
        if timing:
            label += f" ({timing})"
        lines.extend([label, "", NOTES[slide_number], ""])
    return "\n".join(lines).strip() + "\n"


def main():
    if not INPUT_PPTX.exists():
        raise FileNotFoundError(f"Input presentation not found: {INPUT_PPTX}")

    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    presentation = Presentation(str(INPUT_PPTX))

    for slide_number, shape_updates in SLIDE_TEXT.items():
        slide = presentation.slides[slide_number - 1]
        for shape_index, text in shape_updates.items():
            set_shape_text(slide, shape_index, text)

    for slide_number, note_text in NOTES.items():
        slide = presentation.slides[slide_number - 1]
        set_notes(slide, note_text)

    presentation.save(str(OUTPUT_PPTX))
    OUTPUT_SPEECH.write_text(build_speech_markdown(), encoding="utf-8")

    print(f"Saved presentation: {OUTPUT_PPTX}")
    print(f"Saved speech: {OUTPUT_SPEECH}")


if __name__ == "__main__":
    main()
