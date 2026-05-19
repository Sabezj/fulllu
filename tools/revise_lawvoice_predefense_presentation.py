from __future__ import annotations

from pathlib import Path

from pptx import Presentation


SOURCE_PPTX = Path(r"C:\Users\xsanf\Downloads\Telegram Desktop\HSE_LawVoice_PreDefense_Presentation.pptx")
OUTPUT_PPTX = SOURCE_PPTX.with_name("HSE_LawVoice_PreDefense_Presentation_revised_defense.pptx")


def set_shape_text(slide, shape_index: int, text: str) -> None:
    slide.shapes[shape_index].text = text


def revise_presentation() -> None:
    presentation = Presentation(str(SOURCE_PPTX))

    slide = presentation.slides[3]
    set_shape_text(slide, 1, "Goal and scope tightened for defense: backend-centered prototype with measurable evaluation.")
    set_shape_text(
        slide,
        3,
        "Aim\nDesign, implement and evaluate a backend-centered prototype of a real-time retrieval-augmented voice assistant for adolescent legal literacy, with explicit attention to retrieval placement, safety enforcement and measurable backend behavior.",
    )
    set_shape_text(
        slide,
        4,
        "Objectives\nDefine domain and functional requirements.\nCompare voice-assistant and RAG approaches.\nDesign a modular system architecture.\nImplement WebRTC + backend + dialog + retrieval pipeline.\nEvaluate static implementation depth, retrieval metrics, synthetic backend latency, scenarios and risks.\nSpecify observability and scaling limits.",
    )

    slide = presentation.slides[4]
    set_shape_text(slide, 1, "A software-engineering thesis with explicit retrieval, backend and risk measurements.")
    set_shape_text(
        slide,
        5,
        "Evaluation\nStatic codebase metrics\nAPI/test inventory\nProxy RAG metrics: Hit@k, MRR, nDCG\nSynthetic backend latency and throughput\nScenario behavior matrix\nRisk-control assessment and monitoring roadmap",
    )

    slide = presentation.slides[5]
    set_shape_text(slide, 1, "Retrieval is part of model input, not a post-processing layer over model output.")

    slide = presentation.slides[7]
    set_shape_text(slide, 0, "Implementation and measured prototype results")
    set_shape_text(slide, 1, "Static analysis is now complemented by explicit RAG and backend metrics.")
    set_shape_text(slide, 3, "57\nrelevant files")
    set_shape_text(slide, 4, "13\nautomated tests")
    set_shape_text(slide, 5, "0.90\nHit@1")
    set_shape_text(slide, 6, "9.14-16.21 ms\np95 backend routes")

    slide = presentation.slides[9]
    set_shape_text(slide, 1, "Controls reduce risk, but the current prototype cannot claim absolute confidentiality.")

    slide = presentation.slides[10]
    set_shape_text(slide, 1, "Implemented logic supports key scenarios; remaining work is now framed as measurable validation.")
    set_shape_text(
        slide,
        4,
        "Remaining validation\nRun end-to-end voice sessions with credentials and tracing.\nCreate a legally reviewed scenario corpus.\nRun database-backed load tests and measure latency/cost.\nHave experts rate safety, correctness, clarity and empathy.",
    )

    slide = presentation.slides[11]
    set_shape_text(slide, 1, "Personal contribution: integrated architecture, corrected RAG framing, measurable evaluation and defense-ready scope.")
    set_shape_text(
        slide,
        3,
        "What was achieved\nCoherent applied AI system\nRAG placed before generation\nStateful safety policy\nExplicit retrieval and backend metrics\nObservability and scaling plan",
    )
    set_shape_text(
        slide,
        4,
        "Limitations\nNo controlled user study yet\nProxy retrieval benchmark, not final legal benchmark\nProprietary LLM means no full confidentiality claim\nProduction monitoring not yet fully deployed\nExpert review is still required",
    )
    set_shape_text(
        slide,
        5,
        "Future work\nLangfuse request tracing\nPrometheus, Grafana and alerts\nAdversarial evaluation\nDatabase-backed load testing\nExpert and user studies\nHuman-in-the-loop escalation",
    )

    slide = presentation.slides[12]
    set_shape_text(slide, 1, "Best format: what was implemented, what was measured, what can honestly be concluded.")
    set_shape_text(
        slide,
        3,
        "We did\nDesigned and implemented a backend-centered RAG voice-assistant prototype for adolescent legal literacy, with safety routing, retrieval and action planning.",
    )
    set_shape_text(
        slide,
        4,
        "We received\nA measurable prototype description: 57 relevant files, 13 automated tests, proxy retrieval metrics and synthetic backend route benchmarks.",
    )
    set_shape_text(
        slide,
        5,
        "We concluded\nLawVoice is a credible educational prototype if defended as a grounded, risk-reduced backend system rather than as a fully secure or legally authoritative assistant.",
    )

    OUTPUT_PPTX.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(str(OUTPUT_PPTX))


if __name__ == "__main__":
    revise_presentation()
    print(OUTPUT_PPTX)
